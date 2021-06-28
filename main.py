#!/usr/bin/env python3

from argparse import ArgumentParser
from calendar import SUNDAY
from calendar import month_name
from calendar import monthcalendar, setfirstweekday
from datetime import datetime

import jpholiday

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


SLIDE_WIDTH = 33.86
SLIDE_HEIGHT = 19.05

HPOS = 1.69
VPOS = 2.54
WIDTH = 30.56
HEIGHT = 14.94

HEADER_HEIGHT = 0.9
DAY_HEIGHT = 2.86

HEADER = RGBColor(26, 66, 138)
RED = RGBColor(255, 0, 0)
BLUE = RGBColor(0, 145, 218)
BG_WEEKDAY = RGBColor(242, 242, 242)
BG_WEEKEND = RGBColor(217, 217, 217)

DAY_OF_WEEK = ["日", "月", "火", "水", "木", "金", "土"]
FONT_NAME = 'Meiryo'
FONT_SIZE = 12

parser = ArgumentParser()
parser.add_argument("year", type=int, nargs='?', default=datetime.now().year, help="year of calendar target")


def format_title(title):
    title.left = Cm(HPOS)
    title.top = Cm(1.15)
    title.width = Cm(WIDTH)
    title.height = Cm(1.06)
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    title.text_frame.paragraphs[0].font.size = Pt(24)
    title.text_frame.paragraphs[0].font.name = FONT_NAME
    title.text_frame.paragraphs[0].font.color.rgb = HEADER


def format_cell(cell, font_color=False, bg_color=False, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.TOP):
    if bg_color != False:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color
    for p in cell.text_frame.paragraphs:
        if font_color != False:
            p.font.color.rgb = font_color
        p.font.size = Pt(FONT_SIZE)
        p.font.name = FONT_NAME
        p.alignment = align
    cell.vertial_anchor = anchor


def generate_monthly_calendar(slide, cal, holidays):
    shapes = slide.shapes
    width = Cm(WIDTH)
    height = Cm(HEIGHT)
    table = shapes.add_table(len(cal) + 1, 7, Cm(HPOS), Cm(VPOS), width, height).table
    for dn_idx, day_name in enumerate(DAY_OF_WEEK):
        table.cell(0, dn_idx).text = day_name
        format_cell(table.cell(0, dn_idx), False, HEADER, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    table.rows[0].height = Cm(HEADER_HEIGHT)
    for w_idx, week in enumerate(cal):
        for d_idx, day in enumerate(week):
            if day != 0:
                table.cell(w_idx + 1, d_idx).text = str(day)
                if d_idx == 0:
                    format_cell(table.cell(w_idx + 1, d_idx), RED, BG_WEEKEND)
                elif d_idx == 6:
                    format_cell(table.cell(w_idx + 1, d_idx), BLUE, BG_WEEKEND)
                else:
                    format_cell(table.cell(w_idx + 1, d_idx), False, BG_WEEKDAY)
                holiday = list(filter(lambda x: x[0].day == day, holidays))
                if len(holiday):
                    table.cell(w_idx + 1, d_idx).text += "\n" + holiday[0][1]
                    format_cell(table.cell(w_idx + 1, d_idx), RED, BG_WEEKEND)
            else:
                table.cell(w_idx + 1, d_idx).text = '-'
                format_cell(table.cell(w_idx + 1, d_idx), False, BG_WEEKEND)
        table.rows[w_idx + 1].height = Cm(DAY_HEIGHT)


def run(year):
    print('Start generating pptx calendar for year {}'.format(year))
    prs = Presentation()
    prs.slide_width = Cm(SLIDE_WIDTH)
    prs.slide_height = Cm(SLIDE_HEIGHT)
    setfirstweekday(SUNDAY)
    title_only_slide_layout = prs.slide_layouts[5]
    for month in range(1, 13):
        cal = monthcalendar(year, month)
        slide = prs.slides.add_slide(title_only_slide_layout)
        slide.shapes.title.text = '{}/{}'.format(year, month)
        format_title(slide.shapes.title)
        holidays = jpholiday.month_holidays(year, month)
        if len(holidays):
            print('{} has {} holidays'.format(month_name[month], len(holidays)))
            for holiday in holidays:
                print('  {}/{:02d}/{:02d}  {}'.format(holiday[0].year, holiday[0].month, holiday[0].day, holiday[1]))
        generate_monthly_calendar(slide, cal, holidays)
    
    prs.save('{}.pptx'.format(year))
    print('Generation completed. You can find \'{}.pptx\' on current directory.'.format(year))


if __name__ == '__main__':
    args = parser.parse_args()
    run(args.year)

